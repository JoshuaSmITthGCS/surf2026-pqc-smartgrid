#!/usr/bin/env python3
"""
Build the SURF 2026 meeting deck:
"Homomorphic Encryption for Smart-Grid Privacy" — mentor follow-up.

Generates a presenter-ready .pptx that answers Dr. Baza's six questions:
  1. Demand response + HE          4. Which application to implement first
  2. Meter reading + HE            5. Baseline = another HE scheme (Paillier), not RSA
  3. Metrics + who implements      6. Dataset status + action items

Usage:
    pip install python-pptx
    python scripts/build_meeting_deck.py            # -> he_smartgrid_meeting.pptx
    python scripts/build_meeting_deck.py out.pptx   # custom path

Speaker notes are attached to every content slide so you can read them aloud.
"""

from __future__ import annotations

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- theme
NAVY = RGBColor(0x0E, 0x1B, 0x2A)
PANEL = RGBColor(0x16, 0x2A, 0x3F)
ACCENT = RGBColor(0x35, 0xC8, 0xE8)
ACCENT2 = RGBColor(0x7A, 0xE5, 0xB0)
WHITE = RGBColor(0xF2, 0xF6, 0xFA)
MUTED = RGBColor(0xA9, 0xBC, 0xCF)

SW, SH = Inches(13.333), Inches(7.5)  # 16:9
FONT = "Calibri"


def _bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)  # rectangle
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


def _text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6):
    """runs: list of (text, size, color, bold) — one paragraph each."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = txt
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _eyebrow_title(slide, eyebrow, title):
    _box(slide, Inches(0), Inches(0), Inches(0.18), SH, fill=ACCENT)
    _text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.4),
          [(eyebrow.upper(), 13, ACCENT, True)])
    _text(slide, Inches(0.6), Inches(0.85), Inches(12.1), Inches(1.0),
          [(title, 34, WHITE, True)])


def _bullets(slide, items, left=Inches(0.7), top=Inches(2.0),
             width=Inches(12.0), size=18, gap=10):
    """items: list of (text, level) — level 0/1."""
    tb = slide.shapes.add_textbox(left, top, width, SH - top - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.level = lvl
        bullet = "▪  " if lvl == 0 else "–  "
        r = p.add_run()
        r.text = bullet + txt
        r.font.name = FONT
        r.font.size = Pt(size if lvl == 0 else size - 3)
        r.font.color.rgb = WHITE if lvl == 0 else MUTED
        r.font.bold = lvl == 0
    return tb


def _footer(slide, n):
    _text(slide, Inches(0.6), Inches(7.05), Inches(10), Inches(0.3),
          [("Homomorphic Encryption for Smart-Grid Privacy  ·  J. Smith  ·  SURF 2026",
            10, MUTED, False)])
    _text(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.3),
          [(str(n), 10, MUTED, False)], align=PP_ALIGN.RIGHT)


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def new_slide(eyebrow=None, title=None, n=None, footer=True):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    if eyebrow is not None:
        _eyebrow_title(s, eyebrow, title)
    if footer and n is not None:
        _footer(s, n)
    return s


# ============================================================ 1. TITLE
s = new_slide(footer=False)
_box(s, Inches(0), Inches(2.5), SW, Inches(0.06), fill=ACCENT)
_text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.5),
      [("SURF 2026  ·  COLLEGE OF CHARLESTON  ·  MENTOR FOLLOW-UP", 14, ACCENT, True)])
_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6),
      [("Homomorphic Encryption", 46, WHITE, True),
       ("for Smart-Grid Privacy", 46, WHITE, True)])
_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.6),
      [("Researcher:  Joshua Smith", 18, MUTED, False),
       ("Mentor:  Dr. Mohamed Baza", 18, MUTED, False),
       ("Target venue:  IEEE SmartGridComm / CNS", 18, MUTED, False)])
_notes(s, "This session answers the six questions from last meeting: demand "
          "response with HE, meter reading with HE, comparison metrics and who "
          "implements, which application to build first, switching the baseline "
          "to another HE scheme (Paillier) instead of RSA, and the dataset status.")

# ============================================================ 2. AGENDA
s = new_slide("What this meeting answers", "Your Six Questions", 2)
_bullets(s, [
    ("Demand response — how & why is HE useful?", 0),
    ("Meter reading — why HE? privacy only? what is the input data?", 0),
    ("What metrics compare the HE techniques — and who implements them?", 0),
    ("Which HE application do I start implementing first?", 0),
    ("Baseline should be another HE technique, not RSA", 0),
    ("Smart-grid datasets — status, gaps, and next steps", 0),
], top=Inches(2.1), size=20, gap=16)
_notes(s, "I'll take each of your questions in order. Quick headline answers: "
          "start with meter-reading aggregation in BFV; baseline against Paillier "
          "not RSA; ETT data is real but single-transformer and one CSV she sent "
          "is mislabeled.")

# ============================================================ 3. RECAP
s = new_slide("The problem", "Why HE for Smart Grids", 3)
_bullets(s, [
    ("Smart meters stream fine-grained data", 0),
    ("Reveals occupancy, appliances, behavior", 1),
    ("GDPR / CCPA + customer trust demand protection", 1),
    ("Utilities want to outsource compute & storage to the cloud", 0),
    ("…without exposing raw readings to third parties", 1),
], top=Inches(2.0), width=Inches(6.2))
p = _box(s, Inches(7.1), Inches(2.0), Inches(5.6), Inches(4.2), fill=PANEL)
_text(s, Inches(7.4), Inches(2.25), Inches(5.0), Inches(3.8),
      [("THE GAP HE FILLS", 13, ACCENT, True),
       ("Classical encryption (AES) protects data at rest and in transit.", 16, WHITE, False),
       ("But data must be DECRYPTED to compute on it.", 16, MUTED, False),
       ("HE lets you compute directly on ciphertext —", 16, ACCENT2, True),
       ("raw readings are never exposed during processing.", 16, WHITE, False)],
      space_after=12)
_notes(s, "Baseline framing. The novelty is that HE computes on ciphertext, so the "
          "aggregator/cloud never sees a raw reading.")

# ============================================================ 4. METER READING
s = new_slide("Question 2 · Use case A", "Meter Reading with HE", 4)
_bullets(s, [
    ("Meters send kWh per interval for billing & grid operations", 0),
    ("Fine-grained reads leak occupancy / appliances / behavior", 1),
    ("HE value is NOT privacy only:", 0),
    ("Aggregation for load-balancing with no trusted aggregator", 1),
    ("Safe outsourcing of storage/compute to an untrusted cloud", 1),
    ("Time-of-use / tariff billing computed over ciphertext", 1),
    ("Lower utility liability — it never holds raw data", 1),
], top=Inches(2.0), width=Inches(6.6))
p = _box(s, Inches(7.5), Inches(2.0), Inches(5.2), Inches(4.2), fill=PANEL)
_text(s, Inches(7.8), Inches(2.25), Inches(4.6), Inches(3.8),
      [("INPUT DATA FOR HE", 13, ACCENT, True),
       ("Per-meter consumption readings per interval.", 15, WHITE, False),
       ("In the ETT data: load columns HUFL/HULL/MUFL/", 14, MUTED, False),
       ("MULL/LUFL/LULL = the meter \"readings.\"", 14, MUTED, False),
       ("Encode integers (scaled Wh) → BFV / BGV.", 14, ACCENT2, True),
       ("Encode reals directly → CKKS.", 14, ACCENT2, True)],
      space_after=10)
_notes(s, "Answer to 'why HE for meter reading, privacy only?': no — also trustless "
          "aggregation, cloud outsourcing, encrypted billing, reduced liability. The "
          "INPUT is per-meter kWh per interval; I map that to the ETT load columns, "
          "integer-encoded for BFV/BGV and real for CKKS.")

# ============================================================ 5. DEMAND RESPONSE
s = new_slide("Question 1 · Use case B", "Demand Response with HE", 5)
_bullets(s, [
    ("DR = utility signals customers to shift/curtail load at peak", 0),
    ("Must compute: aggregate available flexibility, total reduction", 1),
    ("achieved, and per-customer incentives / settlement", 1),
    ("HE lets the aggregator do this WITHOUT seeing any home's load", 0),
    ("Encrypted sum of per-home curtailment across many meters", 1),
    ("Encrypted comparison of aggregate vs the target", 1),
    ("Fair incentive settlement on encrypted contributions", 1),
    ("DR is MULTI-STEP (baseline→deviation→incentive→settlement)", 0),
    ("Needs multiplicative depth → leveled scheme (BGV/BFV); CKKS forecasts", 1),
], top=Inches(2.0), width=Inches(12.2), size=17, gap=8)
_notes(s, "Answer to 'how & why HE for demand response': the aggregator computes "
          "aggregate flexibility, achieved reduction, and incentive settlement on "
          "ciphertext, never seeing individual consumption or which appliances were "
          "curtailed. Because DR chains several operations it needs depth — that's "
          "the strongest motive for a leveled scheme (BGV/BFV), with CKKS for "
          "real-valued forecasting. Value beyond privacy: enables third-party "
          "aggregators, compliance, and verifiable fair settlement.")

# ============================================================ 6. TRUST MODEL / ACTORS
s = new_slide("Question 3 · Trust model", "Who Does What", 6)


def actor(l, title, sub, color):
    w = Inches(3.4)
    b = _box(s, l, Inches(2.6), w, Inches(1.7), fill=PANEL, line=color)
    _text(s, l + Inches(0.2), Inches(2.75), w - Inches(0.4), Inches(1.4),
          [(title, 18, color, True), (sub, 13, MUTED, False)], space_after=8)
    return l + w


actor(Inches(0.7), "Smart meter", "Encrypts reading under utility public key, "
      "transmits ciphertext. Tight compute → encryption cost & ct size matter.", ACCENT)
actor(Inches(4.95), "Aggregator / cloud", "Untrusted for raw data. Runs the "
      "homomorphic compute (sum / average / statistics) on ciphertext.", ACCENT2)
actor(Inches(9.2), "Utility / key-holder", "Holds secret key. Decrypts only the "
      "final aggregate result — never individual reads.", ACCENT)
_text(s, Inches(0.7), Inches(4.6), Inches(12), Inches(0.6),
      [("Meters  —encrypt→  Aggregator/Cloud  —homomorphic compute→  Key-holder  —decrypt aggregate→  result",
        16, WHITE, True)], align=PP_ALIGN.CENTER)
_bullets(s, [
    ("Meters only encrypt; the utility/cloud never decrypts raw data", 1),
    ("Implication: meter-side encryption cost & ciphertext size are first-class metrics", 1),
], top=Inches(5.4), size=16)
_notes(s, "Answer to 'who implements — smart meters or utility?': both, split by "
          "role. Meters encrypt and send (resource-constrained, so encryption cost "
          "and ciphertext size dominate there). The aggregator/cloud does the "
          "homomorphic computation. The key-holder decrypts only the final "
          "aggregate. This split is why meter-side cost is a headline metric.")

# ============================================================ 7. METRICS COMPUTE
s = new_slide("Question 3 · Metrics", "Computation Time", 7)
_bullets(s, [
    ("Per-operation latency (mean / median / std / p95, ms):", 0),
    ("Key generation", 1),
    ("Encrypt  ·  Decrypt", 1),
    ("Homomorphic add  ·  homomorphic multiply", 1),
    ("Relinearization / rescale (the costly steps)", 1),
    ("End-to-end aggregation latency vs number of meters", 0),
    ("Supported multiplicative depth / remaining noise budget", 0),
    ("Accuracy / error for CKKS (approximate vs plaintext)", 0),
], top=Inches(2.0))
_notes(s, "Reuses the existing BenchmarkRunner harness and fixed CSV schema. The "
          "operations sweep across poly_modulus_degree. CKKS adds an accuracy "
          "metric because it's approximate.")

# ============================================================ 8. METRICS COMMS
s = new_slide("Question 3 · Metrics", "Communication Overhead", 8)
_bullets(s, [
    ("Ciphertext size and ciphertext EXPANSION ratio (bytes / reading)", 0),
    ("Key sizes: public key, relinearization keys, Galois keys", 0),
    ("Total bytes per aggregation round (uplink + downlink)", 0),
    ("Memory / storage footprint at the aggregator", 0),
    ("Scalability: latency & bytes vs #meters and vs poly_modulus_degree", 0),
], top=Inches(2.0), width=Inches(7.0))
p = _box(s, Inches(7.9), Inches(2.0), Inches(4.8), Inches(3.6), fill=PANEL)
_text(s, Inches(8.2), Inches(2.25), Inches(4.2), Inches(3.2),
      [("WHY IT MATTERS", 13, ACCENT, True),
       ("Meters are bandwidth- and power-constrained.", 15, WHITE, False),
       ("RLWE ciphertexts are large; expansion ratio is", 15, MUTED, False),
       ("often the real bottleneck, not CPU time.", 15, MUTED, False),
       ("Paillier vs RLWE differ sharply here.", 15, ACCENT2, True)],
      space_after=10)
_notes(s, "Communication is often the binding constraint on the meter side — RLWE "
          "ciphertexts are large. Reporting the expansion ratio makes the "
          "Paillier-vs-RLWE trade-off concrete.")

# ============================================================ 9. BASELINE
s = new_slide("Question 5 · Baseline change", "Baseline = Paillier, not RSA", 9)
_bullets(s, [
    ("Per your feedback: compare HE against HE, not RSA/AES", 0),
    ("Paillier = additively homomorphic (PHE), the classic prior-", 0),
    ("generation smart-metering baseline (Erkin & Tsudik, ACNS 2012)", 1),
    ("Head-to-head: Paillier (PHE)  vs  BFV / CKKS / BGV (RLWE)", 0),
    ("Plus intra-RLWE comparison: BFV vs CKKS vs BGV", 0),
    ("AES-GCM / RSA demoted to a one-line context mention only", 0),
], top=Inches(2.0), width=Inches(12.2), size=18, gap=12)
_notes(s, "Directly addresses 'baseline should be another HE technique not RSA.' "
          "Paillier is additively homomorphic — perfect for the encrypted-sum "
          "meter-reading use case and the standard prior baseline in the smart-"
          "metering literature. AES/RSA stay only as one line of context.")

# ============================================================ 10. COMPARISON TABLE
s = new_slide("Question 5 · Scheme landscape", "Paillier vs BFV vs CKKS vs BGV", 10)
rows = [
    ["Scheme", "Data type", "Exactness", "Homomorphism", "Best fit", "Library"],
    ["Paillier", "Integers", "Exact", "Additive only (PHE)", "Baseline encrypted sums", "phe / custom"],
    ["BFV", "Integers", "Exact", "Add + mult, fixed mod", "Exact meter aggregation", "TenSEAL"],
    ["CKKS", "Real / complex", "Approx.", "Add + mult, rescaling", "Stats / anomaly / ML", "TenSEAL"],
    ["BGV", "Integers", "Exact", "Add + mult, leveled", "Deep DR pipelines", "SEAL"],
]
nrows, ncols = len(rows), len(rows[0])
tbl_shape = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(2.1),
                               Inches(12.1), Inches(3.6))
table = tbl_shape.table
widths = [1.7, 1.8, 1.5, 2.7, 2.7, 1.7]
for c, w in enumerate(widths):
    table.columns[c].width = Inches(w)
for r in range(nrows):
    for c in range(ncols):
        cell = table.cell(r, c)
        cell.text = rows[r][c]
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.name = FONT
        run.font.size = Pt(13 if r else 13)
        run.font.bold = (r == 0) or (c == 0)
        run.font.color.rgb = NAVY if r == 0 else WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else (PANEL if r % 2 else NAVY)
_notes(s, "One-glance comparison. Note Paillier is additive-only — fine for sums but "
          "can't multiply, which is exactly why DR's multi-step pipeline needs RLWE. "
          "BGV needs Microsoft SEAL; TenSEAL doesn't expose it.")

# ============================================================ 11. WHERE TO START
s = new_slide("Question 4 · First implementation", "Start: BFV Meter Aggregation", 11)
_bullets(s, [
    ("Recommendation: meter-reading / demand aggregation in BFV first", 0),
    ("Simplest correct circuit — additions only, low depth", 1),
    ("BFV already runs in TenSEAL (no new library needed)", 1),
    ("Gives a clean first benchmark milestone vs the Paillier baseline", 1),
    ("Then CKKS for encrypted average / variance (anomaly stats)", 0),
    ("BGV (via SEAL) later for the deeper demand-response pipeline", 0),
], top=Inches(2.0), width=Inches(12.2), size=18, gap=12)
_notes(s, "Answer to 'which application to start — meter reading?': yes. BFV exact "
          "integer summation for meter-reading aggregation is the cleanest first "
          "milestone and is already running. CKKS and BGV follow.")

# ============================================================ 12. DATASETS
s = new_slide("Question 6 · Data status", "Datasets — Honest Picture", 12)
_bullets(s, [
    ("ETTh1 / ETTh2 (hourly), ETTm1 (15-min): real Electricity", 0),
    ("Transformer Temperature data — 6 load features + oil temp", 1),
    ("Good for load aggregation & transformer-anomaly framing", 1),
    ("BUT one transformer, not many homes → multi-meter must be", 1),
    ("SIMULATED (columns / time-windows / resampled segments)", 1),
    ("Customer_Behaviour.csv is the \"Social Network Ads\" dataset", 0),
    ("(Age / EstimatedSalary / Purchased) — NOT grid telemetry", 1),
    ("Flag as mislabeled; at most a toy encrypted-classification demo", 1),
], top=Inches(2.0), width=Inches(12.2), size=17, gap=8)
_notes(s, "Be honest: ETT is real and usable but single-transformer, so 'aggregate "
          "across many meters' has to be simulated. The Customer_Behaviour file is "
          "actually the Social Network Ads dataset — not smart-grid data at all.")

# ============================================================ 13. DATA ACTIONS
s = new_slide("Question 6 · Next steps", "Dataset Action Items", 13)
_bullets(s, [
    ("Find a true multi-HOUSEHOLD smart-meter dataset. Candidates:", 0),
    ("UMass Smart*  ·  Pecan Street / Dataport", 1),
    ("UK-DALE  ·  CER Irish Smart Metering  ·  Low Carbon London", 1),
    ("Email bluhmja@g.cofc.edu to follow up on available datasets", 0),
    ("Ask \"adawa\" for Adam's email and request the dataset", 0),
    ("Decide: real multi-home data vs documented simulation from ETT", 0),
], top=Inches(2.0), width=Inches(12.2), size=18, gap=12)
_notes(s, "Concrete follow-ups: name candidate public datasets, email "
          "bluhmja@g.cofc.edu, and get Adam's email via adawa. Until then, document "
          "the simulated-meter approach on ETT.")

# ============================================================ 14. PLAN
s = new_slide("Plan", "Next Two Weeks", 14)
_bullets(s, [
    ("Implement BFV meter-reading aggregation end-to-end", 0),
    ("Add Paillier baseline for the same encrypted-sum task", 0),
    ("Benchmark: compute time + ciphertext expansion, Paillier vs BFV", 0),
    ("Wire results through the existing BenchmarkRunner CSV schema", 0),
    ("Resolve dataset: secure multi-home data or finalize ETT simulation", 0),
    ("Then: CKKS averages/variance; scope BGV via SEAL", 0),
], top=Inches(2.0), width=Inches(12.2), size=18, gap=12)
_notes(s, "A focused two-week milestone set tied to the decisions above.")

# ============================================================ 15. QUESTIONS
s = new_slide("Discussion", "Open Questions for You", 15)
_bullets(s, [
    ("Paillier as the HE baseline — agreed, or prefer another PHE?", 0),
    ("Meter-reading aggregation as the first build — agreed?", 0),
    ("Is simulating many meters from ETT acceptable for now?", 0),
    ("Should I prioritize getting real multi-home data before benchmarking?", 0),
    ("Is BGV (via SEAL) in-scope this summer, or stretch goal?", 0),
], top=Inches(2.0), width=Inches(12.2), size=19, gap=14)
_notes(s, "End on decisions I need from you so we leave aligned.")

# ---------------------------------------------------------------- save
out = sys.argv[1] if len(sys.argv) > 1 else "he_smartgrid_meeting.pptx"
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides._sldIdLst)} slides)")
